from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import tempfile
import os
import subprocess

# Helper: add page number to slide
def add_page_number(slide, page_num):
    left = Inches(8.5) - Inches(1.2)
    top = Inches(7.5) - Inches(0.5)
    width = Inches(1)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    # Ensure only one paragraph and set just the numeric slide number
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{page_num}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.RIGHT

# Helper: export slide as image (requires PowerPoint or LibreOffice for best results)
import hashlib
import shutil

def _format_markdown_text(text_frame, markdown_text):
    """Format markdown-like text in PowerPoint text frame with basic formatting."""
    import re
    lines = markdown_text.split('\n')
    
    # Clear existing content
    text_frame.clear()
    
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            # Add empty paragraph for spacing
            if i > 0:
                text_frame.add_paragraph()
            continue
            
        # Add new paragraph for each line
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Handle different markdown elements
        if line.startswith('# '):
            # Header 1
            p.text = line[2:].strip()
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
        elif line.startswith('## '):
            # Header 2
            p.text = line[3:].strip()
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
        elif line.startswith('### '):
            # Header 3
            p.text = line[4:].strip()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
        elif line.startswith('- ') or line.startswith('* '):
            # Bullet point
            _add_formatted_text(p, line[2:].strip(), 12)
        elif re.match(r'^\d+\. ', line):
            # Numbered list
            text = re.sub(r'^\d+\. ', '', line).strip()
            _add_formatted_text(p, text, 12)
        else:
            # Regular text with potential inline formatting
            _add_formatted_text(p, line, 12)

def _add_formatted_text(paragraph, text, font_size):
    """Add text with inline bold and italic formatting to a paragraph."""
    import re
    
    # Split text by both **bold** and *italic* markers
    # Use a more complex regex to handle both bold and italic
    parts = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*)', text)
    
    # Clear paragraph and start fresh
    paragraph.clear()
    
    for i, part in enumerate(parts):
        if part.startswith('**') and part.endswith('**') and len(part) > 4:
            # Bold text
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = RGBColor(30, 30, 30)
        elif part.startswith('*') and part.endswith('*') and len(part) > 2 and not part.startswith('**'):
            # Italic text (single asterisks, not double)
            run = paragraph.add_run()
            run.text = part[1:-1]
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = RGBColor(30, 30, 30)
        elif part:
            # Regular text
            run = paragraph.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = RGBColor(30, 30, 30)

def get_slide_images(input_pptx):
    # Hash the pptx file to create a unique temp folder
    with open(input_pptx, 'rb') as f:
        pptx_hash = hashlib.md5(f.read()).hexdigest()
    base = os.path.basename(input_pptx)
    name = os.path.splitext(base)[0]
    base_dir = os.path.join('temp', f'{name}-{pptx_hash}')
    img_dir = os.path.join(base_dir, 'slides')
    os.makedirs(img_dir, exist_ok=True)

    # Reuse if already exported
    first_img = os.path.join(img_dir, 'slide_output-1.png')
    if os.path.exists(first_img):
        return img_dir

    # Locate LibreOffice
    soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice_path is None:
        raise RuntimeError("LibreOffice (soffice) is not installed or not in PATH. Cannot export slides.")

    # 1) Convert PPTX to PDF into base_dir
    try:
        subprocess.run([
            soffice_path, "--headless", "--convert-to", "pdf", "--outdir", base_dir, input_pptx
        ], check=True)
    except subprocess.CalledProcessError:
        raise RuntimeError("LibreOffice failed to convert PPTX to PDF.")

    pdf_path = os.path.join(base_dir, f"{name}.pdf")
    if not os.path.exists(pdf_path):
        # Some LO builds may name output differently; attempt to find the first .pdf in base_dir
        candidates = [f for f in os.listdir(base_dir) if f.lower().endswith('.pdf')]
        if candidates:
            pdf_path = os.path.join(base_dir, candidates[0])
        else:
            raise RuntimeError("PDF not found after LibreOffice conversion.")

    # 2) Convert PDF to PNG images using pdftoppm
    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm is None:
        raise RuntimeError("pdftoppm (poppler-utils) is not installed. Please install it to export images.")

    prefix = os.path.join(img_dir, 'slide_output')
    try:
        subprocess.run([pdftoppm, "-png", pdf_path, prefix], check=True)
    except subprocess.CalledProcessError:
        raise RuntimeError("pdftoppm failed to convert PDF to images.")

    return img_dir

def export_slide_as_image(prs, slide_idx, tmpdir, input_pptx=None):
    # Use real slide image if available
    if input_pptx:
        img_dir = get_slide_images(input_pptx)
        img_path = os.path.join(img_dir, f'slide_output-{slide_idx+1}.png')
        if os.path.exists(img_path):
            return img_path
    # fallback: placeholder
    img_path = os.path.join(tmpdir, f"slide_{slide_idx+1}.png")
    img = Image.new('RGB', (1280, 720), color=(240, 240, 240))
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 80)
    except:
        font = ImageFont.load_default()
    draw.text((img.width//2-100, img.height//2-40), f"Slide {slide_idx+1}", fill=(100,100,100), font=font)
    img.save(img_path)
    return img_path

# Main processing function
def process_presentation(input_pptx, output_pptx, notes_path: str | None = None):
    prs = Presentation(input_pptx)
    tmpdir = tempfile.mkdtemp()
    orig_slide_count = len(prs.slides)
    slide_indices = list(range(orig_slide_count))
    page_num = 1
    # Load notes if provided
    notes_by_slide = [""] * orig_slide_count
    if notes_path:
        try:
            from utils.notes_parser import parse_notes
            notes_by_slide = parse_notes(notes_path, orig_slide_count)
        except Exception as e:
            print(f"Warning: failed to parse notes file: {e}")
    for offset, slide_idx in enumerate(slide_indices):
        slide = prs.slides[slide_idx + offset]
        add_page_number(slide, page_num)
        # Export slide as image (real image if possible)
        img_path = export_slide_as_image(prs, slide_idx, tmpdir, input_pptx=input_pptx)
        # Duplicate slide after current
        blank_slide_layout = prs.slide_layouts[6]  # blank
        new_slide = prs.slides.add_slide(blank_slide_layout)
        # Move new slide to correct position
        prs.slides._sldIdLst.insert(slide_idx + offset + 1, prs.slides._sldIdLst[-1])
        # Add image to new slide (upper-right, preserve aspect ratio at ~30% slide width)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        margin = int(slide_width * 0.05)
        try:
            with Image.open(img_path) as im:
                img_w, img_h = im.size
        except Exception:
            img_w, img_h = (1280, 720)
        # target width is 30% of slide width (50% bigger than previous 20%)
        target_w = int(slide_width * 0.30)
        target_h = int(target_w * img_h / img_w)
        # ensure height fits within slide minus margins
        max_h = int(slide_height - 2 * margin)
        if target_h > max_h and max_h > 0:
            target_h = max_h
            target_w = int(target_h * img_w / img_h)
        left = int(slide_width - target_w - margin)
        top = int(slide_height * 0.05)
        new_slide.shapes.add_picture(img_path, left, top, width=target_w, height=target_h)
        # Add notes textbox filling the remaining left area if notes exist
        note_text = notes_by_slide[slide_idx] if slide_idx < len(notes_by_slide) else ""
        if note_text:
            gap = int(slide_width * 0.02)
            notes_left = int(margin)
            # Make sure notes don't go above the screenshot
            notes_top = max(int(slide_height * 0.05), top + int(target_h * 0.1))
            notes_right = left - gap
            notes_width = max(0, notes_right - notes_left)
            notes_bottom = int(slide_height - margin)
            notes_height = max(0, notes_bottom - notes_top)
            if notes_width > 0 and notes_height > 0:
                tb = new_slide.shapes.add_textbox(notes_left, notes_top, notes_width, notes_height)
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = None
                # Format Markdown-like text in PowerPoint
                _format_markdown_text(tf, note_text)
        page_num += 1
    prs.save(output_pptx)
    print(f"Saved: {output_pptx}")
